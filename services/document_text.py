# services/document_text.py
import io
import json
import base64
import mimetypes
from pathlib import Path

import requests
from aiogram.types import Message

from config import YANDEX_API_KEY, YANDEX_FOLDER_ID, YANDEX_OCR_URL


# ===== ВСПОМОГАТЕЛЬНОЕ: ЗАГРУЗКА ФАЙЛА ИЗ TELEGRAM =====

async def _download_file_bytes(message: Message) -> bytes:
    """
    Скачиваем файл/фото из Telegram и возвращаем raw bytes.
    """
    file_obj = None

    if message.document:
        file_obj = await message.bot.get_file(message.document.file_id)
    elif message.photo:
        # берем самое большое фото
        file_obj = await message.bot.get_file(message.photo[-1].file_id)

    if not file_obj:
        return b""

    file_bytes = await message.bot.download_file(
        file_obj.file_path,
        destination=io.BytesIO()
    )
    file_bytes.seek(0)
    return file_bytes.read()


# ===== OCR ДЛЯ ИЗОБРАЖЕНИЙ / СКАНОВ =====

def _call_ocr(image_bytes: bytes, mime_type: str = "JPEG") -> str:
    """
    Вызов Yandex OCR и получение текста (fullText + строки).
    """
    if not image_bytes:
        return ""

    content_b64 = base64.b64encode(image_bytes).decode("utf-8")

    data = {
        "mimeType": mime_type,          # JPEG / PNG / PDF
        "languageCodes": ["ru", "en"],  # можно ["*"] для автоопределения
        "model": "page",
        "content": content_b64,
    }

    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Api-Key {YANDEX_API_KEY}",
        "x-folder-id": YANDEX_FOLDER_ID,
        "x-data-logging-enabled": "true",
    }

    resp = requests.post(
        url=YANDEX_OCR_URL,
        headers=headers,
        data=json.dumps(data),
    )
    print("OCR RESPONSE:", resp.status_code, resp.text)
    resp.raise_for_status()
    result = resp.json()

    text_annotation = result.get("result", {}).get("textAnnotation", {})
    full_text = text_annotation.get("fullText")
    if full_text:
        return full_text.strip()

    lines = []
    for block in text_annotation.get("blocks", []):
        for line in block.get("lines", []):
            txt = line.get("text", "").strip()
            if txt:
                lines.append(txt)

    return "\n".join(lines)


# ===== ЛОКАЛЬНЫЙ ПАРСИНГ ТЕКСТОВЫХ ФАЙЛОВ =====

def _extract_text_from_pdf(file_bytes: bytes) -> str:
    """
    Пробуем достать текст из PDF без OCR (если PDF текстовый).
    """
    from pypdf import PdfReader  # pip install pypdf

    buf = io.BytesIO(file_bytes)
    reader = PdfReader(buf)
    pieces = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        if txt:
            pieces.append(txt)
    return "\n".join(pieces).strip()


def _extract_text_from_docx(file_bytes: bytes) -> str:
    """
    DOCX: вытаскиваем текст абзацев.
    """
    from docx import Document  # pip install python-docx

    buf = io.BytesIO(file_bytes)
    doc = Document(buf)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras).strip()


def _extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    PPTX: собираем текст из всех слайдов.
    """
    from pptx import Presentation  # pip install python-pptx

    buf = io.BytesIO(file_bytes)
    prs = Presentation(buf)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    lines.append(txt)
    return "\n".join(lines).strip()


def _extract_text_from_txt(file_bytes: bytes) -> str:
    """
    Простой текстовый файл (txt, md, csv и т.п.).
    """
    try:
        return file_bytes.decode("utf-8").strip()
    except UnicodeDecodeError:
        # запасной вариант
        return file_bytes.decode("cp1251", errors="ignore").strip()


def _guess_extension(filename: str | None) -> str:
    if not filename:
        return ""
    return Path(filename).suffix.lower()  # ".pdf", ".docx" и т.п.


# ===== ОСНОВНАЯ ФУНКЦИЯ ДЛЯ БОТА =====

async def extract_text(message: Message) -> str:
    """
    Универсальная функция:
    - если фото → OCR
    - если документ → в зависимости от расширения:
        * .txt/.md/.csv/... → читаем как текст
        * .pdf → сначала пытаемся вытащить текст, если пусто — OCR
        * .docx → python-docx
        * .pptx → python-pptx
      если формат неизвестен → пробуем как текст, если не вышло — OCR как изображение/PDF
    """
    # 1. Фото → сразу OCR
    if message.photo:
        image_bytes = await _download_file_bytes(message)
        text = _call_ocr(image_bytes, mime_type="JPEG")
        return text or "Фотография с текстом, но OCR не вернул результат."

    # 2. Документ
    if message.document:
        file_bytes = await _download_file_bytes(message)
        filename = message.document.file_name or ""
        ext = _guess_extension(filename)

        # Простые текстовые форматы
        if ext in {".txt", ".md", ".csv", ".log"}:
            text = _extract_text_from_txt(file_bytes)
            return text or f"Документ {filename}, но текст не удалось прочитать."

        # DOCX
        if ext == ".docx":
            try:
                text = _extract_text_from_docx(file_bytes)
                if text:
                    return text
            except Exception as e:
                print("DOCX parse error:", e)

        # PPTX
        if ext == ".pptx":
            try:
                text = _extract_text_from_pptx(file_bytes)
                if text:
                    return text
            except Exception as e:
                print("PPTX parse error:", e)

        # PDF: сначала «текстовый» парсинг, потом OCR
        if ext == ".pdf":
            try:
                text = _extract_text_from_pdf(file_bytes)
                if text:
                    return text
            except Exception as e:
                print("PDF parse error:", e)

            # если парсинг не дал результата — пробуем OCR по PDF
            ocr_text = _call_ocr(file_bytes, mime_type="PDF")
            return ocr_text or f"PDF {filename}, но ни текст, ни OCR не дали результата."

        # Неизвестный формат:
        # 1) пробуем прочитать как текст
        text = _extract_text_from_txt(file_bytes)
        if text:
            return text

        # 2) fallback — OCR как изображение (на случай сканов в непонятных форматах)
        ocr_text = _call_ocr(file_bytes, mime_type="JPEG")
        return ocr_text or f"Документ: {filename}, но формат не распознан."

    # 3. Ничего не передали
    return "Неизвестный документ"
